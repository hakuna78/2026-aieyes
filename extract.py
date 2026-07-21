import os
import fitz
import traceback
from pptx import Presentation
import comtypes.client

input_dir = r"C:\Users\samsung\Downloads\2026-aieyes-main\2026-aieyes-main\public\ideathon"
output_dir = r"C:\Users\samsung\Downloads\2026-aieyes-main\2026-aieyes-main\public\ideathon_thumbs"

files = ["1조.pptx.pptx", "2조.pdf", "3조.pptx", "4조.pdf", "5조 발표.pdf"]

try:
    powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
    powerpoint.Visible = 1
except Exception as e:
    powerpoint = None

titles = {}

for f in files:
    file_path = os.path.join(input_dir, f)
    thumb_path = os.path.join(output_dir, f + ".png")
    
    title = ""
    
    if f.lower().endswith('.pdf'):
        try:
            doc = fitz.open(file_path)
            page = doc[0]
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
            pix.save(thumb_path)
            
            text = page.get_text("text").strip()
            lines = [l.strip() for l in text.split('\n') if l.strip()]
            if lines:
                title = lines[0]
            doc.close()
        except Exception as e:
            pass
            
    elif f.lower().endswith('.pptx'):
        try:
            prs = Presentation(file_path)
            for shape in prs.slides[0].shapes:
                if not shape.has_text_frame:
                    continue
                if shape.text.strip():
                    title = shape.text.strip()
                    break
            
            if powerpoint:
                presentation = powerpoint.Presentations.Open(file_path, WithWindow=False)
                presentation.Slides[1].Export(thumb_path, "PNG", 1280, 720)
                presentation.Close()
        except Exception as e:
            pass
            
    titles[f] = title

with open("titles.txt", "w", encoding="utf-8") as f:
    for k, v in titles.items():
        f.write(f"{k}:::{v}\n")

if powerpoint:
    powerpoint.Quit()
